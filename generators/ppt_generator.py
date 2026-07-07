"""
ThinkForge AI
PowerPoint Generator

Creates a professional PowerPoint presentation
from ProjectState.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

from core.state import ProjectState


class PPTGenerator:

    def __init__(self):

        self.output_dir = Path(
            "outputs/ppt"
        )

        self.output_dir.mkdir(
            parents=True,
            exist_ok=True
        )

    ########################################################

    def add_title_slide(

        self,

        prs: Presentation,

        title: str,

        subtitle: str

    ):

        layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = title

        slide.placeholders[1].text = subtitle

    ########################################################

    def add_bullet_slide(

        self,

        prs,

        title,

        bullets

    ):

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = title

        body = slide.placeholders[1].text_frame

        body.clear()

        for item in bullets:

            p = body.add_paragraph()

            p.text = str(item)

            p.level = 0

            p.font.size = Pt(22)

    ########################################################

    def add_text_slide(

        self,

        prs,

        title,

        text

    ):

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = title

        body = slide.placeholders[1]

        body.text = text

    ########################################################

    def add_budget_slide(

        self,

        prs,

        budget

    ):

        layout = prs.slide_layouts[5]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = "Estimated Budget"

        rows = len(
            budget["items"]
        ) + 2

        cols = 2

        table = slide.shapes.add_table(

            rows,

            cols,

            Inches(0.8),

            Inches(1.2),

            Inches(8),

            Inches(3.5)

        ).table

        table.cell(
            0,
            0
        ).text = "Item"

        table.cell(
            0,
            1
        ).text = "Cost"

        row = 1

        for item in budget["items"]:

            table.cell(
                row,
                0
            ).text = item["name"]

            table.cell(
                row,
                1
            ).text = item["cost"]

            row += 1

        table.cell(
            row,
            0
        ).text = "TOTAL"

        table.cell(
            row,
            1
        ).text = budget["total"]
        ########################################################

    def add_timeline_slide(

        self,

        prs,

        timeline

    ):

        bullets = []

        for phase in timeline:

            bullets.append(

                f"{phase['phase']} : {phase['duration']}"

            )

        self.add_bullet_slide(

            prs,

            "Implementation Timeline",

            bullets

        )

    ########################################################

    def add_risk_slide(

        self,

        prs,

        risks

    ):

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = "Risk Analysis"

        tf = slide.placeholders[1].text_frame

        tf.clear()

        for risk in risks:

            p = tf.add_paragraph()

            p.text = (

                f"Risk: {risk['risk']}\n"

                f"Mitigation: {risk['mitigation']}"

            )

            p.level = 0

            p.font.size = Pt(18)

    ########################################################

    def add_kpi_slide(

        self,

        prs,

        kpis

    ):

        self.add_bullet_slide(

            prs,

            "Success KPIs",

            kpis

        )

    ########################################################

    def add_architecture_slide(

        self,

        prs,

        architecture

    ):

        bullets = architecture.get(

            "components",

            []

        )

        self.add_bullet_slide(

            prs,

            architecture.get(

                "title",

                "Architecture"

            ),

            bullets

        )

    ########################################################

    def add_expert_slide(

        self,

        prs,

        experts

    ):

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = "Expert Council"

        tf = slide.placeholders[1].text_frame

        tf.clear()

        for expert in experts:

            p = tf.add_paragraph()

            p.text = (

                f"{expert['name']}"

                f" ({expert['role']})"

            )

            p.level = 0

            p.font.size = Pt(20)

    ########################################################

    def add_debate_slide(

        self,

        prs,

        debate

    ):

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = "Expert Consensus"

        tf = slide.placeholders[1].text_frame

        tf.clear()

        p = tf.add_paragraph()

        p.text = debate.get(

            "final_consensus",

            ""

        )

        p.font.size = Pt(22)

    ########################################################

    def add_report_slide(

        self,

        prs,

        report

    ):

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = "Technical Report"

        body = slide.placeholders[1]

        if len(report) > 1000:

            report = report[:1000] + "..."

        body.text = report

    ########################################################

    def add_confidence_slide(

        self,

        prs,

        confidence

    ):

        layout = prs.slide_layouts[5]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = "AI Confidence"

        tb = slide.shapes.add_textbox(

            Inches(2),

            Inches(2),

            Inches(5),

            Inches(2)

        )

        tf = tb.text_frame

        p = tf.paragraphs[0]

        p.text = (

            f"{round(confidence*100,2)}%"

        )

        p.font.size = Pt(34)  
        ########################################################
    # Main Generator
    ########################################################

    def generate(
        self,
        state: ProjectState
    ) -> Path:

        prs = Presentation()

        ####################################################
        # Title
        ####################################################

        self.add_title_slide(
            prs,
            "ThinkForge AI",
            "Turn Problems Into Deployable Solutions"
        )

        ####################################################
        # Problem
        ####################################################

        self.add_text_slide(
            prs,
            "Problem Statement",
            state.problem
        )

        ####################################################
        # Objective
        ####################################################

        self.add_text_slide(
            prs,
            "Project Objective",
            state.objective
        )

        ####################################################
        # Executive Summary
        ####################################################

        self.add_text_slide(
            prs,
            "Executive Summary",
            state.executive_summary
        )

        ####################################################
        # Domain
        ####################################################

        self.add_text_slide(
            prs,
            "Business Domain",
            state.domain
        )

        ####################################################
        # Stakeholders
        ####################################################

        self.add_bullet_slide(
            prs,
            "Stakeholders",
            state.stakeholders
        )

        ####################################################
        # Constraints
        ####################################################

        self.add_bullet_slide(
            prs,
            "Constraints",
            state.constraints
        )

        ####################################################
        # Assumptions
        ####################################################

        self.add_bullet_slide(
            prs,
            "Assumptions",
            state.assumptions
        )

        ####################################################
        # Experts
        ####################################################

        self.add_expert_slide(
            prs,
            state.experts
        )

        ####################################################
        # Debate
        ####################################################

        self.add_debate_slide(
            prs,
            state.debate
        )

        ####################################################
        # Architecture
        ####################################################

        self.add_architecture_slide(
            prs,
            state.architecture
        )

        ####################################################
        # Budget
        ####################################################

        self.add_budget_slide(
            prs,
            state.budget
        )

        ####################################################
        # Timeline
        ####################################################

        self.add_timeline_slide(
            prs,
            state.timeline
        )

        ####################################################
        # Risks
        ####################################################

        self.add_risk_slide(
            prs,
            state.risks
        )

        ####################################################
        # KPIs
        ####################################################

        self.add_kpi_slide(
            prs,
            state.kpis
        )

        ####################################################
        # Roadmap
        ####################################################

        if isinstance(
            state.implementation_plan,
            list
        ):

            self.add_bullet_slide(
                prs,
                "Implementation Roadmap",
                state.implementation_plan
            )

        ####################################################
        # Technical Report
        ####################################################

        self.add_report_slide(
            prs,
            state.technical_report
        )

        ####################################################
        # Confidence
        ####################################################

        self.add_confidence_slide(
            prs,
            state.confidence
        )

        ####################################################
        # Thank You
        ####################################################

        self.add_title_slide(
            prs,
            "Thank You",
            "Generated automatically by ThinkForge AI"
        )

        ####################################################
        # Save
        ####################################################

        output_file = self.output_dir / "ThinkForge_Presentation.pptx"

        prs.save(output_file)

        return output_file